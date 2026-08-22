# -*- coding: utf-8 -*-
"""IR 덱 디자인 계약 + 헬퍼 라이브러리 (지침 22 §3).

콘텐츠·디자인 분리(지침 22 §2): 텍스트는 deck.md(SSOT)에, 디자인은 여기 THEMES에.
디자인은 *여러 후보*로 탐색한다 — THEMES에 테마를 추가하고 build.py --theme <id>로 렌더해 고른다.

계약 구조(색 3단계·타입 6단계·강조 예산제)는 테마와 무관하게 유지한다(지침 20 "이름=역할").
후보 간에 바뀌는 것은 *값·성격*(다크↔라이트, 코랄↔딥블루)이지 *구조*가 아니다.
슬라이드 코드는 토큰 이름만 부른다 — raw hex/pt 금지. 토큰은 호출 시점에 활성 테마로 해석된다.
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


def _rgb(hexstr):
    h = hexstr.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ============================================================
# 타입 스케일 6단계 (§3.2) — 구조 상수. 테마 무관(디스플레이 크기는 예외)
# ============================================================
T_HEAD = 26    # 헤드라인(완결 주장문)
T_CARD = 14    # 카드·섹션 제목
T_SUB  = 12    # 소제목/키커
T_BODY = 10.5  # 본문
T_CAP  = 9.5   # 캡션
T_FOOT = 9     # 각주·푸터

SW, SH = 13.333, 7.5   # 16:9
MX = 0.7               # 좌우 마진
FOOT_LINE = 7.05       # 하단 요소 y+h ≤ 이 값 (§7 겹침 가드)
WHITE = _rgb("FFFFFF")

# ============================================================
# THEMES — 디자인 후보 레지스트리 (§3, P3). 새 후보는 여기 한 줄
# ============================================================
# 각 테마 = 색 3단계(TXT_HI/MID/LOW) + 베이스/강조 + 배경 + 서체.
# 값·성격만 다르게, 구조(키 집합)는 동일하게 유지한다.
THEMES = {
    # A. 다크 네이비 프리미엄 (icons-invest v4 정본)
    "A_darknavy": {
        "name": "다크 네이비 프리미엄",
        "BG":     _rgb("1A1F4B"),  # 슬라이드 기본 배경
        "BASE":   _rgb("2A3170"),  # 톤온톤 보조
        "ACCENT": _rgb("FF5C7A"),  # ★ 강조 1색 — 장당 1~2포인트 예산(§3.3)
        "TXT_HI":  _rgb("FFFFFF"),
        "TXT_MID": _rgb("D5DBF0"),
        "TXT_LOW": _rgb("9AA2C8"),  # WCAG AA 4.5:1 검증 대상
        "MUTED":  _rgb("6B728F"),
        "FONT":   "Pretendard",
    },
    # B. 라이트 벤토 미니멀 (밝은 배경·차분한 강조)
    "B_lightbento": {
        "name": "라이트 벤토 미니멀",
        "BG":     _rgb("FFFFFF"),
        "BASE":   _rgb("F4F6FC"),
        "ACCENT": _rgb("3B4CB8"),
        "TXT_HI":  _rgb("1A1F4B"),   # 밝은 배경 → 잉크가 핵심색
        "TXT_MID": _rgb("3A4166"),
        "TXT_LOW": _rgb("6B728F"),
        "MUTED":  _rgb("9AA2C8"),
        "FONT":   "Pretendard",
    },
    # C. 딥 코퍼레이트 (절제된 기업 톤 — HYBE 공개자료 계열)
    "C_deepcorp": {
        "name": "딥 코퍼레이트",
        "BG":     _rgb("14161C"),
        "BASE":   _rgb("23262F"),
        "ACCENT": _rgb("C8A24B"),
        "TXT_HI":  _rgb("FFFFFF"),
        "TXT_MID": _rgb("C9CCD6"),
        "TXT_LOW": _rgb("8A8F9C"),
        "MUTED":  _rgb("6B6F7A"),
        "FONT":   "Pretendard",
    },
}

# 활성 테마 토큰 (apply_theme가 채운다). 이름 = 역할.
BG = BASE = ACCENT = TXT_HI = TXT_MID = TXT_LOW = MUTED = None
HEAD = BODY = None
ACTIVE = None


def apply_theme(theme_id):
    """활성 테마의 토큰을 모듈 전역에 바인딩. 헬퍼는 호출 시점에 이 전역을 읽는다."""
    global BG, BASE, ACCENT, TXT_HI, TXT_MID, TXT_LOW, MUTED, HEAD, BODY, ACTIVE
    t = THEMES[theme_id]
    BG, BASE, ACCENT = t["BG"], t["BASE"], t["ACCENT"]
    TXT_HI, TXT_MID, TXT_LOW, MUTED = t["TXT_HI"], t["TXT_MID"], t["TXT_LOW"], t["MUTED"]
    HEAD = BODY = t["FONT"]
    ACTIVE = theme_id
    return t


apply_theme("A_darknavy")  # import 시 기본 테마 — build.py에서 --theme로 교체

# ============================================================
# HELPERS — 슬라이드 코드가 부르는 원자적 프리미티브 (토큰 late-bind)
# ============================================================
def slide(prs, bg=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill = BG if bg is None else bg
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = fill; r.line.fill.background(); r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def pic(s, path, l, t, w, h):
    return s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))


def rect(s, l, t, w, h, fill=None, line=None, lw=1.0, round=False, radius=0.08):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if round:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def txt(s, l, t, w, h, runs, size=T_BODY, bold=False, color=None,
        align=PP_ALIGN.LEFT, font=None, anchor=MSO_ANCHOR.TOP, sp=1.0, wrap=True):
    color = TXT_MID if color is None else color
    font = BODY if font is None else font
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(runs, str): runs = [(runs, {})]
    def is_line_list(x): return isinstance(x, list) and x and isinstance(x[0], list)
    lines = runs if is_line_list(runs) else [runs]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = sp
        if isinstance(ln, tuple): ln = [ln]
        for t_, o in ln:
            r = p.add_run(); r.text = t_
            r.font.size = Pt(o.get('size', size)); r.font.bold = o.get('bold', bold)
            r.font.name = o.get('font', font); r.font.color.rgb = o.get('color', color)
            if o.get('italic'): r.font.italic = True
    return tb


def title(s, head, kicker=None, color=None):
    """헤드라인 = 완결 주장문(§3.4). 명사구 금지."""
    if kicker:
        txt(s, MX, 0.45, 12, 0.3, [(kicker, {'size': T_SUB, 'bold': True, 'color': ACCENT, 'font': HEAD})])
    txt(s, MX, 0.72, SW - 2 * MX, 1.15, head, size=T_HEAD, bold=True,
        color=(TXT_HI if color is None else color), font=HEAD, sp=1.02)


def foot(s, text_):
    """각주·출처 — 전장 동일 위치·9pt·TXT_LOW (§3.1)."""
    txt(s, MX, SH - 0.42, SW - 2 * MX, 0.3, [(text_, {'size': T_FOOT, 'color': TXT_LOW})])
