# -*- coding: utf-8 -*-
"""IR 덱 빌드 오케스트레이터 (지침 22 P3~P5).

콘텐츠·디자인 분리: 텍스트는 deck.md(SSOT), 디자인은 contract.THEMES에서 고른다.

사용:
  python3 build.py                     # 활성/선택 테마로 전체 덱 빌드
  python3 build.py --theme B_lightbento # 특정 테마로 빌드
  python3 build.py --candidates        # 모든 테마의 대표 슬라이드를 각각 렌더(후보 비교, P3)

필요: pip install python-pptx ; LibreOffice(soffice) PATH; 렌더 서체(Pretendard) 설치.
"""
import os
import sys
import subprocess as sp
from pptx import Presentation
from pptx.util import Inches
import contract as C
from contract import *  # noqa: 헬퍼 + 활성 테마 토큰

HERE = os.path.dirname(os.path.abspath(__file__))
DATA = os.path.join(HERE, "_data")
CHART_GEN = os.path.join(HERE, "charts.py")  # xlsx 정본 → PNG (있을 때만)


# ── 슬라이드 정의 — deck.md(SSOT) 텍스트를 반영. 1슬라이드=1함수, 원자적 ──
# 텍스트를 바꿀 때는 deck.md에서 먼저 고치고 여기 반영한다(지침 22 §2 불변규율 1).
def s01_cover(prs):
    s = slide(prs)
    txt(s, 0, 1.95, SW, 2.0, [("BRAND", {'size': 104, 'bold': True, 'color': TXT_HI, 'font': HEAD})],
        align=PP_ALIGN.CENTER)
    txt(s, 0, 4.25, SW, 0.45, [("한 줄 태그라인.", {'size': 16, 'color': TXT_MID})], align=PP_ALIGN.CENTER)
    txt(s, 0, SH - 0.62, SW, 0.3,
        [("㈜회사  ·  2026.MM  ·  Strictly Private & Confidential", {'size': T_FOOT, 'color': TXT_LOW})],
        align=PP_ALIGN.CENTER)


def s02_market(prs):
    s = slide(prs)
    title(s, [("시장은 이렇게 움직이고 있습니다 — 완결 주장문으로",
               {'size': T_HEAD, 'bold': True, 'color': TXT_HI, 'font': HEAD})], kicker="시장 기회")
    # 강조 예산제(§3.3): ACCENT는 장당 1~2포인트만
    rect(s, MX, 2.2, 3.6, 2.4, fill=BASE, round=True)
    txt(s, MX + 0.3, 2.5, 3.0, 0.4, [("64%", {'size': 40, 'bold': True, 'color': ACCENT, 'font': HEAD})])
    txt(s, MX + 0.3, 3.4, 3.0, 1.0, [("코어팬 30%가 팬덤 소비의 64%", {'size': T_BODY, 'color': TXT_MID})])
    foot(s, "출처: (기관, 2026)")


def s03_close(prs):
    s = slide(prs)
    txt(s, 0, 3.0, SW, 1.2, [("함께 만들 다음 장.", {'size': 48, 'bold': True, 'color': TXT_HI, 'font': HEAD})],
        align=PP_ALIGN.CENTER)


# 대표 슬라이드(후보 비교용) vs 전체 덱. 실제 덱은 SLIDES에 계속 추가.
REPRESENTATIVE = [s01_cover, s02_market, s03_close]
SLIDES = [s01_cover, s02_market, s03_close]


def build(theme_id, slide_fns, out_name):
    C.apply_theme(theme_id)
    prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    for fn in slide_fns:
        fn(prs)
    out = os.path.join(HERE, out_name)
    prs.save(out)
    _geometry_check(prs, out_name)
    _render_pdf(out)
    return out


def _geometry_check(prs, label):
    """하단 겹침 가드(§7): 하단 요소 y+h ≤ 7.05. 전면 배경은 제외."""
    warns = []
    for i, s in enumerate(prs.slides, 1):
        for shp in s.shapes:
            try:
                top = shp.top / 914400
                bottom = (shp.top + shp.height) / 914400  # EMU→inch
                is_full_bleed = top < 0.05 and bottom >= SH - 0.05
                if not is_full_bleed and bottom > FOOT_LINE + 0.35:
                    warns.append(f"  {label} S{i:02d}: {shp.shape_type} bottom={bottom:.2f} > {FOOT_LINE}")
            except (TypeError, AttributeError):
                pass
    print("[warn] 하단 겹침 후보:\n" + "\n".join(warns) if warns else f"[ok] {label} geometry 통과")


def _render_pdf(pptx_path):
    try:
        sp.run(["soffice", "--headless", "--convert-to", "pdf", "--outdir", HERE, pptx_path],
               check=True, capture_output=True)
        print(f"[ok] PDF 동기화: {os.path.basename(pptx_path)[:-5]}.pdf")
    except (sp.CalledProcessError, FileNotFoundError) as e:
        print(f"[skip] soffice 렌더 불가: {e}. 서체 임베드는 pdffonts로 검증할 것.")


# ── P0: 차트 데이터 정본(xlsx) → PNG 자동 재생성 ──────────────
if os.path.exists(CHART_GEN):
    sp.run(["python3", CHART_GEN], check=True)

if __name__ == "__main__":
    args = sys.argv[1:]
    if "--candidates" in args:
        # P3: 모든 테마의 대표 슬라이드를 각각 렌더 → 나란히 비교 후 선택
        print("[P3] 디자인 후보 렌더 — 대표 슬라이드로 비교:")
        for tid, t in THEMES.items():
            build(tid, REPRESENTATIVE, f"candidate_{tid}.pptx")
            print(f"     → candidate_{tid} ({t['name']})")
        print("고른 테마를 build.py --theme <id> 또는 contract.apply_theme 기본값으로 고정(§3).")
    else:
        theme = "A_darknavy"
        if "--theme" in args:
            theme = args[args.index("--theme") + 1]
        out = build(theme, SLIDES, "deck.pptx")
        print(f"[ok] {out}  ({len(SLIDES)} slides · theme={theme})")
